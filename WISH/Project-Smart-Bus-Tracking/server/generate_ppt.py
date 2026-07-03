"""
Generate a polished, editorial-style presentation (.pptx) for the rural
bus-tracking project. Inspired by the RuralSync layout: warm cream background,
bold display titles, line-art accent illustrations in the corners.

Dependencies (already on the system):
    pip install python-pptx Pillow

Run:
    python3 generate_ppt.py
    open Project_Presentation.pptx
"""

from io import BytesIO
import os
import math

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from PIL import Image, ImageDraw

# ============================================================
# Theme
# ============================================================
BG_CREAM   = RGBColor(0xF5, 0xF0, 0xE6)   # warm off-white
INK_DARK   = RGBColor(0x14, 0x18, 0x24)   # near-black titles
INK_BODY   = RGBColor(0x2A, 0x2E, 0x3A)   # body
INK_MUTE   = RGBColor(0x6B, 0x70, 0x7D)   # captions
ACCENT     = RGBColor(0xC2, 0x41, 0x0C)   # burnt orange
ACCENT_SOFT= RGBColor(0xE9, 0xB7, 0x9A)   # soft peach

# PIL versions (RGB tuples)
PIL_ACCENT      = (0xC2, 0x41, 0x0C)
PIL_ACCENT_SOFT = (0xE9, 0xB7, 0x9A)
PIL_BG          = (0xF5, 0xF0, 0xE6, 0)   # transparent canvas

TITLE_FONT = "Arial Black"   # bold display fallback present on macOS / Win
BODY_FONT  = "Helvetica"
SERIF_FONT = "Georgia"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ============================================================
# Line-art generators (Pillow) -> PNG bytes for embedding
# ============================================================

def _new_canvas(w=1200, h=1200):
    img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    return img, ImageDraw.Draw(img)


def _png_bytes(img):
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def art_skyline(color=PIL_ACCENT):
    """A simple city-skyline silhouette of rectangles + a couple of spires."""
    img, d = _new_canvas(1600, 900)
    y = 880
    rects = [
        (40,   500, 180, y),
        (180,  430, 260, y),
        (260,  560, 350, y),
        (350,  380, 460, y),  # tall
        (440,  470, 540, y),
        (540,  300, 640, y),  # tallest
        (640,  480, 740, y),
        (740,  410, 850, y),
        (850,  520, 960, y),
        (960,  360, 1070, y),
        (1070, 470, 1180, y),
        (1180, 540, 1290, y),
        (1290, 410, 1400, y),
        (1400, 500, 1520, y),
    ]
    for x0, y0, x1, y1 in rects:
        d.rectangle([x0, y0, x1, y1], outline=color, width=4)
        # a few windows
        rows = max(1, (y1 - y0) // 60)
        for r in range(rows):
            wy = y0 + 20 + r * 60
            if wy + 25 < y1 - 20:
                d.rectangle([x0 + 12, wy, x0 + 32, wy + 25], outline=color, width=2)
                d.rectangle([x0 + (x1 - x0) // 2 + 6, wy, x0 + (x1 - x0) // 2 + 26, wy + 25], outline=color, width=2)

    # antenna on the tallest
    d.line([(590, 300), (590, 230)], fill=color, width=4)
    d.line([(585, 240), (595, 240)], fill=color, width=3)
    # ground line
    d.line([(0, y), (1600, y)], fill=color, width=3)
    return _png_bytes(img)


def art_bus(color=PIL_ACCENT):
    img, d = _new_canvas(1400, 700)
    # body
    d.rounded_rectangle([60, 120, 1240, 520], radius=60, outline=color, width=6)
    # roof line
    d.line([(60, 200), (1240, 200)], fill=color, width=4)
    # windows
    for i in range(6):
        x0 = 130 + i * 170
        d.rounded_rectangle([x0, 230, x0 + 130, 360], radius=12, outline=color, width=4)
    # door
    d.rounded_rectangle([1120, 230, 1210, 480], radius=8, outline=color, width=4)
    d.line([(1165, 230), (1165, 480)], fill=color, width=3)
    # wheels
    d.ellipse([200, 460, 360, 620], outline=color, width=6)
    d.ellipse([240, 500, 320, 580], outline=color, width=4)
    d.ellipse([920, 460, 1080, 620], outline=color, width=6)
    d.ellipse([960, 500, 1040, 580], outline=color, width=4)
    # headlight
    d.ellipse([60, 380, 110, 430], outline=color, width=4)
    # signal waves
    for r in (40, 80, 120):
        d.arc([300 - r, 30 - r, 300 + r, 30 + r], 200, 340, fill=color, width=4)
    for r in (40, 80, 120):
        d.arc([900 - r, 30 - r, 900 + r, 30 + r], 200, 340, fill=color, width=4)
    return _png_bytes(img)


def art_pin(color=PIL_ACCENT):
    img, d = _new_canvas(600, 800)
    # teardrop
    d.ellipse([100, 80, 500, 480], outline=color, width=10)
    # bottom point - draw two lines for the spike
    d.polygon([(300, 720), (180, 420), (420, 420)], outline=color, fill=None)
    d.line([(180, 420), (300, 720)], fill=color, width=10)
    d.line([(420, 420), (300, 720)], fill=color, width=10)
    # inner dot
    d.ellipse([240, 220, 360, 340], outline=color, width=8)
    # rays
    for ang in range(0, 360, 30):
        x0 = 300 + int(360 * math.cos(math.radians(ang)))
        y0 = 280 + int(360 * math.sin(math.radians(ang)))
        x1 = 300 + int(420 * math.cos(math.radians(ang)))
        y1 = 280 + int(420 * math.sin(math.radians(ang)))
        d.line([(x0, y0), (x1, y1)], fill=color, width=5)
    return _png_bytes(img)


def art_lightbulb(color=PIL_ACCENT):
    img, d = _new_canvas(700, 900)
    # bulb
    d.ellipse([120, 80, 580, 540], outline=color, width=8)
    # base
    d.rounded_rectangle([240, 540, 460, 640], radius=12, outline=color, width=6)
    d.line([(240, 580), (460, 580)], fill=color, width=4)
    d.line([(240, 610), (460, 610)], fill=color, width=4)
    # screw cap
    d.polygon([(290, 640), (410, 640), (380, 740), (320, 740)], outline=color)
    d.line([(290, 640), (320, 740)], fill=color, width=6)
    d.line([(410, 640), (380, 740)], fill=color, width=6)
    d.line([(320, 740), (380, 740)], fill=color, width=6)
    # filament squiggle
    d.line([(260, 330), (310, 280), (360, 360), (410, 280), (450, 340)], fill=color, width=5)
    # rays
    cx, cy, r0, r1 = 350, 300, 280, 360
    for ang in (-90, -60, -30, 0, 30, 60, 90, 120, 150, 180, 210, 240, 270):
        x0 = cx + int(r0 * math.cos(math.radians(ang)))
        y0 = cy + int(r0 * math.sin(math.radians(ang)))
        x1 = cx + int(r1 * math.cos(math.radians(ang)))
        y1 = cy + int(r1 * math.sin(math.radians(ang)))
        d.line([(x0, y0), (x1, y1)], fill=color, width=5)
    return _png_bytes(img)


def art_gear(color=PIL_ACCENT):
    img, d = _new_canvas(800, 800)
    cx, cy = 400, 400
    rO, rI = 320, 230
    # teeth
    for i in range(12):
        a = i * 30
        x0 = cx + int(rI * math.cos(math.radians(a - 8)))
        y0 = cy + int(rI * math.sin(math.radians(a - 8)))
        x1 = cx + int(rO * math.cos(math.radians(a - 8)))
        y1 = cy + int(rO * math.sin(math.radians(a - 8)))
        x2 = cx + int(rO * math.cos(math.radians(a + 8)))
        y2 = cy + int(rO * math.sin(math.radians(a + 8)))
        x3 = cx + int(rI * math.cos(math.radians(a + 8)))
        y3 = cy + int(rI * math.sin(math.radians(a + 8)))
        d.polygon([(x0, y0), (x1, y1), (x2, y2), (x3, y3)], outline=color)
    d.ellipse([cx - rI, cy - rI, cx + rI, cy + rI], outline=color, width=8)
    d.ellipse([cx - 80, cy - 80, cx + 80, cy + 80], outline=color, width=8)
    return _png_bytes(img)


def art_route(color=PIL_ACCENT):
    img, d = _new_canvas(1400, 800)
    # dotted curving route
    pts = [(60, 600), (260, 400), (520, 520), (780, 280), (1040, 420), (1340, 220)]
    # connect via small line segments to draw curve
    last = pts[0]
    steps = 30
    seg_pts = []
    for i in range(len(pts) - 1):
        for t in range(steps + 1):
            tt = t / steps
            # cubic-ish interp via midpoint
            x = (1 - tt) * pts[i][0] + tt * pts[i + 1][0]
            y = (1 - tt) * pts[i][1] + tt * pts[i + 1][1]
            seg_pts.append((x, y))
    # draw dashed
    for i in range(0, len(seg_pts) - 1, 4):
        x0, y0 = seg_pts[i]
        x1, y1 = seg_pts[min(i + 2, len(seg_pts) - 1)]
        d.line([(x0, y0), (x1, y1)], fill=color, width=6)
    # stop markers
    for p in pts:
        d.ellipse([p[0] - 22, p[1] - 22, p[0] + 22, p[1] + 22], outline=color, width=6)
        d.ellipse([p[0] - 8, p[1] - 8, p[0] + 8, p[1] + 8], fill=color)
    return _png_bytes(img)


def art_busstop(color=PIL_ACCENT):
    img, d = _new_canvas(900, 900)
    # roof
    d.line([(60, 240), (840, 240)], fill=color, width=8)
    d.line([(60, 240), (60, 280)], fill=color, width=6)
    d.line([(840, 240), (840, 280)], fill=color, width=6)
    d.line([(60, 280), (840, 280)], fill=color, width=4)
    # poles
    d.line([(140, 280), (140, 800)], fill=color, width=8)
    d.line([(760, 280), (760, 800)], fill=color, width=8)
    # bench
    d.rectangle([200, 620, 700, 660], outline=color, width=6)
    d.line([(220, 660), (220, 760)], fill=color, width=6)
    d.line([(680, 660), (680, 760)], fill=color, width=6)
    d.line([(450, 660), (450, 760)], fill=color, width=4)
    # sign
    d.rounded_rectangle([360, 90, 540, 220], radius=20, outline=color, width=6)
    d.text((420, 130), "BUS", fill=color)
    # ground
    d.line([(0, 800), (900, 800)], fill=color, width=4)
    return _png_bytes(img)


def art_dataflow(color=PIL_ACCENT, dark=(0x14, 0x18, 0x24)):
    """Slide-sized architecture diagram: Bus -> Ingestion -> Engine -> DB -> Passenger."""
    W, H = 2400, 900
    img = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)

    boxes = [
        ( 60, "Bus\nGPS Hardware"),
        (520, "Person 1\nIngestion + DB"),
        (980, "Person 2\nETA Engine"),
        (1440, "Database\nPositions / ETAs"),
        (1900, "Person 3\nSMS / IVR / Dash"),
    ]
    box_w, box_h = 380, 280
    y0 = 300
    for x, label in boxes:
        d.rounded_rectangle([x, y0, x + box_w, y0 + box_h], radius=24, outline=color, width=6)
        # label rendering done in PPT instead — but draw a tiny header bar
        d.rectangle([x, y0, x + box_w, y0 + 8], fill=color)
    # arrows
    for i in range(len(boxes) - 1):
        x_start = boxes[i][0] + box_w + 10
        x_end   = boxes[i + 1][0] - 10
        yy = y0 + box_h // 2
        d.line([(x_start, yy), (x_end, yy)], fill=dark, width=5)
        # arrowhead
        d.polygon([(x_end, yy), (x_end - 20, yy - 12), (x_end - 20, yy + 12)], fill=dark)
    return _png_bytes(img), boxes, box_w, box_h, y0


# ============================================================
# Slide content
# ============================================================
SLIDES = [
    {
        "type": "title",
        "title": "Rural Bus Tracking &\nPassenger Information System",
        "subtitle": "Low-bandwidth GPS tracking, ETA prediction, multilingual alerts",
        "footer": "WISH Project   ·   IIT Bombay   ·   Summer 2026",
    },
    {
        "type": "section",
        "number": "01",
        "title": "Problem & Motivation",
    },
    {
        "type": "twocol",
        "title": "The Problem",
        "lead": "Rural passengers have no reliable way to know when a bus will arrive — and existing trackers assume continuous 4G and a smartphone.",
        "bullets": [
            "Patchy or absent mobile data along rural routes",
            "Feature phones dominate; smartphone apps don't reach the audience",
            "Multiple regional languages required, not just English / Hindi",
            "Operators have no live view of route adherence or delays",
        ],
        "art": "busstop",
    },
    {
        "type": "twocol",
        "title": "Project Goals",
        "lead": "Build a system that meets passengers where they are — feature phone, SMS, IVR — and gives operators real visibility.",
        "bullets": [
            "Track buses with low-cost hardware over intermittent networks",
            "Predict accurate ETAs even during signal blackouts",
            "Reach every passenger via SMS / IVR — no smartphone required",
            "Alert passengers automatically when a bus diverts from its route",
        ],
        "art": "pin",
    },
    {
        "type": "section",
        "number": "02",
        "title": "System Architecture",
    },
    {
        "type": "dataflow",
        "title": "End-to-End Data Flow",
        "labels": [
            ("Bus", "GPS Hardware"),
            ("Person 1", "Ingestion + DB"),
            ("Person 2", "ETA Engine"),
            ("Database", "Positions / ETAs"),
            ("Person 3", "SMS · IVR · Dash"),
        ],
        "caption": "Each layer owns a clear responsibility; the database is the contract between them.",
    },
    {
        "type": "twocol",
        "title": "Three-Layer Architecture",
        "lead": "Responsibilities split across three contributors, glued together by a shared PostgreSQL + PostGIS database.",
        "bullets": [
            "Bus hardware emits a GPS ping every 2 – 3 minutes",
            "Person 1 — ingests, validates, deduplicates packets",
            "Person 2 — derives position, ETA, status, divert alerts",
            "Person 3 — delivers SMS, IVR, dashboard, stop-side display",
        ],
        "art": "gear",
    },
    {
        "type": "section",
        "number": "03",
        "title": "Person 1 — Ingestion & Database",
    },
    {
        "type": "twocol",
        "title": "Responsibilities (P1)",
        "lead": "The data plumbing layer — converts raw bus pings into a clean stream the engine can trust.",
        "bullets": [
            "Receive GPS packets via HTTP, fall back to SMS gateway when offline",
            "Validate, deduplicate, and clean incoming pings",
            "Maintain canonical route data: stops + road polyline",
            "Expose a clean stream of positions to the ETA engine",
        ],
        "art": "route",
    },
    {
        "type": "twocol",
        "title": "Route Data — Where the Polyline Comes From",
        "lead": "A practical fallback chain so every route has geometry, even when authoritative data is missing.",
        "bullets": [
            "OpenStreetMap + OSRM — automated polyline fetch (default)",
            "Manual digitisation in QGIS / Google My Maps for missing roads",
            "One test drive per route — most accurate, also seeds travel times",
            "GTFS shapes.txt where available; Google Directions as paid fallback",
        ],
        "art": "pin",
    },
    {
        "type": "section",
        "number": "04",
        "title": "Person 2 — ETA & Prediction Engine",
    },
    {
        "type": "twocol",
        "title": "The Four Jobs of the Engine",
        "lead": "Turns raw lat/lon into something a passenger can actually act on.",
        "bullets": [
            "ETA calculation — distance ÷ speed, refined by learned segment times",
            "Network-loss prediction — dead-reckoning when pings stop",
            "Historical learning — per-segment, per-time-of-day averages",
            "Route adherence — 50 m geofence corridor with auto-alert on divert",
        ],
        "art": "gear",
    },
    {
        "type": "states",
        "title": "Prediction State Machine",
        "lead": "The status flag drives wording on every passenger-facing surface.",
        "states": [
            ("LIVE",       "Fresh GPS ping in the last few minutes"),
            ("ESTIMATED",  "Projected from last known speed during short blackout"),
            ("LAST_KNOWN", "Blackout too long — freeze position, mark uncertain"),
            ("OFF_ROUTE",  "Three consecutive packets outside the corridor"),
        ],
    },
    {
        "type": "twocol",
        "title": "Cold-Start Strategy",
        "lead": "On day one the historical table is empty. A four-step fallback chain means the engine never fails to produce an ETA.",
        "bullets": [
            "Hardcoded defaults — ~30 km/h highway, ~25 mixed, ~15 town",
            "Seed from the published bus schedule (stop-to-stop times)",
            "Seed from OSRM driving-time estimates per segment",
            "Optional test drive per route for the highest-fidelity seed",
            "Fallback chain: bucket → segment → route → default",
        ],
        "art": "lightbulb",
    },
    {
        "type": "twocol",
        "title": "What the Engine Writes to the DB",
        "lead": "A single contract: Person 3 reads these rows, the engine guarantees them.",
        "bullets": [
            "Bus position (lat, lon, timestamp)",
            "Next stop and ETAs for all upcoming stops",
            "Status flag — LIVE / ESTIMATED / LAST_KNOWN / OFF_ROUTE",
            "Confidence score 0–1 — drives wording (\"12 min\" vs \"~12 min\")",
        ],
        "art": "route",
    },
    {
        "type": "section",
        "number": "05",
        "title": "Person 3 — Passenger Interface",
    },
    {
        "type": "twocol",
        "title": "Reaching Every Passenger",
        "lead": "Smartphone-free by design. SMS and IVR work on any handset; the dashboard and display board cover operators and waiting passengers.",
        "bullets": [
            "SMS query handler: text a bus number, get an ETA back",
            "IVR call flow: speaks the ETA in the chosen regional language",
            "Authority dashboard: live map, alerts, adherence view",
            "Stop-side display board: next bus + ETA in local language",
            "Maintains the query log that drives divert auto-alerts",
        ],
        "art": "busstop",
    },
    {
        "type": "section",
        "number": "06",
        "title": "Divert Alert Pipeline",
    },
    {
        "type": "twocol",
        "title": "End-to-End Divert Flow",
        "lead": "Fully automated — no operator action required to notify affected passengers.",
        "bullets": [
            "Engine detects 3 consecutive off-corridor packets → OFF_ROUTE",
            "Writes alert row to DB, flashes banner on authority dashboard",
            "Looks up recent query log for that bus (last N hours)",
            "Auto-sends SMS / IVR to every recent querier in their language",
            "Operator notified simultaneously — no manual step",
        ],
        "art": "pin",
    },
    {
        "type": "section",
        "number": "07",
        "title": "Tech Stack & Roadmap",
    },
    {
        "type": "twocol",
        "title": "Tech Stack",
        "lead": "Open-source where possible, paid services only where they earn it.",
        "bullets": [
            "Backend — Python (FastAPI / Flask)",
            "Database — PostgreSQL + PostGIS for geospatial queries",
            "Routing — self-hosted OSRM over OpenStreetMap data",
            "Telephony — SMS gateway + IVR provider (Exotel / Twilio)",
            "Hardware — low-cost GPS module + 2G modem on each bus",
        ],
        "art": "gear",
    },
    {
        "type": "twocol",
        "title": "Status & Pilot Plan",
        "lead": "Prototype works end-to-end on simulated data; next step is hardware-in-the-loop.",
        "bullets": [
            "Phase 1 — bench test with simulator (done)",
            "Phase 2 — single bus, single route, schedule-seeded ETAs",
            "Phase 3 — 5 routes, dashboard + SMS live, ground-truth collection",
            "Phase 4 — historical learning active, accuracy review at 4 weeks",
        ],
        "art": "route",
    },
    {
        "type": "title",
        "title": "Thank You",
        "subtitle": "Questions & Discussion",
        "footer": "WISH Project   ·   IIT Bombay",
    },
]


# ============================================================
# Pillow-art cache
# ============================================================
ART_BUILDERS = {
    "skyline":    art_skyline,
    "bus":        art_bus,
    "pin":        art_pin,
    "lightbulb":  art_lightbulb,
    "gear":       art_gear,
    "route":      art_route,
    "busstop":    art_busstop,
}
_art_cache = {}

def get_art(name):
    if name not in _art_cache:
        _art_cache[name] = ART_BUILDERS[name]().getvalue()
    return BytesIO(_art_cache[name])


# ============================================================
# Slide primitives
# ============================================================
def fill_background(slide, color=BG_CREAM):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid(); rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    # send to back
    spTree = rect._element.getparent()
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)


def add_text(slide, left, top, width, height, text, *, font=BODY_FONT,
             size=18, bold=False, italic=False, color=INK_BODY, align=None,
             line_spacing=1.15, space_after=0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align == "center": p.alignment = PP_ALIGN.CENTER
        elif align == "right": p.alignment = PP_ALIGN.RIGHT
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_thin_rule(slide, left, top, width, color=ACCENT, thick=Pt(2)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thick)
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_art(slide, name, left, top, width, height=None):
    img = get_art(name)
    if height is None:
        slide.shapes.add_picture(img, left, top, width=width)
    else:
        slide.shapes.add_picture(img, left, top, width=width, height=height)


def slide_number_badge(slide, n, total):
    add_text(slide, Inches(12.4), Inches(7.05), Inches(0.7), Inches(0.3),
             f"{n:02d} / {total:02d}", font=BODY_FONT, size=10, color=INK_MUTE, align="right")


# ============================================================
# Slide renderers
# ============================================================
def render_title(slide, data):
    fill_background(slide)
    # corner art
    add_art(slide, "skyline", Inches(-0.6), Inches(4.3), width=Inches(8.5))
    add_art(slide, "pin",     Inches(11.0), Inches(0.4), width=Inches(2.2))

    # decorative serif "tag"
    add_text(slide, Inches(0.9), Inches(1.1), Inches(6), Inches(0.4),
             "WISH  ·  IIT BOMBAY  ·  2026", font=BODY_FONT, size=12, bold=True,
             color=ACCENT)
    add_thin_rule(slide, Inches(0.9), Inches(1.55), Inches(0.6))

    # big bold title
    add_text(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(3.0),
             data["title"], font=TITLE_FONT, size=58, bold=True,
             color=INK_DARK, line_spacing=1.05)

    add_text(slide, Inches(0.9), Inches(5.1), Inches(11.5), Inches(0.6),
             data["subtitle"], font=SERIF_FONT, size=20, italic=True,
             color=INK_BODY)

    add_text(slide, Inches(0.9), Inches(7.0), Inches(11.5), Inches(0.4),
             data["footer"], font=BODY_FONT, size=11, color=INK_MUTE)


def render_section(slide, data):
    fill_background(slide)
    # big section number on the left
    add_text(slide, Inches(0.7), Inches(1.5), Inches(4), Inches(4),
             data["number"], font=TITLE_FONT, size=220, bold=True,
             color=ACCENT_SOFT, line_spacing=1.0)
    # section title
    add_text(slide, Inches(5.0), Inches(3.0), Inches(8.0), Inches(1.5),
             "SECTION", font=BODY_FONT, size=14, bold=True, color=ACCENT)
    add_thin_rule(slide, Inches(5.0), Inches(3.45), Inches(0.8))
    add_text(slide, Inches(5.0), Inches(3.7), Inches(8.0), Inches(2.0),
             data["title"], font=TITLE_FONT, size=48, bold=True,
             color=INK_DARK, line_spacing=1.05)


def render_twocol(slide, data):
    fill_background(slide)

    # top eyebrow + rule
    add_text(slide, Inches(0.7), Inches(0.55), Inches(6), Inches(0.3),
             "RURALSYNC", font=BODY_FONT, size=11, bold=True, color=ACCENT)
    add_thin_rule(slide, Inches(0.7), Inches(0.9), Inches(0.6))

    # title
    add_text(slide, Inches(0.7), Inches(1.05), Inches(12), Inches(1.2),
             data["title"], font=TITLE_FONT, size=36, bold=True,
             color=INK_DARK, line_spacing=1.1)

    # lead
    add_text(slide, Inches(0.7), Inches(2.3), Inches(8.0), Inches(1.2),
             data["lead"], font=SERIF_FONT, size=17, italic=True,
             color=INK_BODY, line_spacing=1.3)

    # bullets
    box = slide.shapes.add_textbox(Inches(0.7), Inches(3.7), Inches(8.0), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(data["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        run1 = p.add_run()
        run1.text = "—  "
        run1.font.name = BODY_FONT
        run1.font.size = Pt(18)
        run1.font.bold = True
        run1.font.color.rgb = ACCENT
        run2 = p.add_run()
        run2.text = b
        run2.font.name = BODY_FONT
        run2.font.size = Pt(18)
        run2.font.color.rgb = INK_BODY

    # right-side accent art
    art_name = data.get("art", "gear")
    add_art(slide, art_name, Inches(9.4), Inches(2.5), width=Inches(3.4))


def render_dataflow(slide, data):
    fill_background(slide)
    add_text(slide, Inches(0.7), Inches(0.55), Inches(6), Inches(0.3),
             "ARCHITECTURE", font=BODY_FONT, size=11, bold=True, color=ACCENT)
    add_thin_rule(slide, Inches(0.7), Inches(0.9), Inches(0.6))
    add_text(slide, Inches(0.7), Inches(1.05), Inches(12), Inches(1.0),
             data["title"], font=TITLE_FONT, size=36, bold=True, color=INK_DARK)

    # five boxes
    labels = data["labels"]
    n = len(labels)
    box_w = 2.0
    box_h = 1.5
    total_w = n * box_w + (n - 1) * 0.35
    left0 = (13.333 - total_w) / 2
    y_box = 3.0

    for i, (top_line, sub_line) in enumerate(labels):
        x = left0 + i * (box_w + 0.35)
        # rounded card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(y_box),
                                       Inches(box_w), Inches(box_h))
        card.fill.solid(); card.fill.fore_color.rgb = BG_CREAM
        card.line.color.rgb = ACCENT
        card.line.width = Pt(1.5)
        card.shadow.inherit = False
        # top bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(x), Inches(y_box),
                                     Inches(box_w), Inches(0.12))
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        # text
        add_text(slide, Inches(x), Inches(y_box + 0.3), Inches(box_w), Inches(0.4),
                 top_line, font=TITLE_FONT, size=14, bold=True, color=INK_DARK,
                 align="center")
        add_text(slide, Inches(x), Inches(y_box + 0.75), Inches(box_w), Inches(0.6),
                 sub_line, font=BODY_FONT, size=11, color=INK_MUTE, align="center")

        # arrow to next
        if i < n - 1:
            ax0 = x + box_w + 0.04
            ax1 = x + box_w + 0.31
            ay  = y_box + box_h / 2
            line = slide.shapes.add_connector(1, Inches(ax0), Inches(ay),
                                              Inches(ax1), Inches(ay))
            line.line.color.rgb = INK_DARK
            line.line.width = Pt(1.5)
            # tiny arrow triangle
            tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                         Inches(ax1 - 0.08), Inches(ay - 0.06),
                                         Inches(0.12), Inches(0.12))
            tri.rotation = 30
            tri.fill.solid(); tri.fill.fore_color.rgb = INK_DARK
            tri.line.fill.background()

    # caption
    add_text(slide, Inches(0.7), Inches(5.6), Inches(12), Inches(0.6),
             data["caption"], font=SERIF_FONT, size=15, italic=True,
             color=INK_MUTE, align="center")

    # decorative art
    add_art(slide, "route", Inches(0.6), Inches(6.0), width=Inches(5.0))


def render_states(slide, data):
    fill_background(slide)
    add_text(slide, Inches(0.7), Inches(0.55), Inches(6), Inches(0.3),
             "ENGINE", font=BODY_FONT, size=11, bold=True, color=ACCENT)
    add_thin_rule(slide, Inches(0.7), Inches(0.9), Inches(0.6))
    add_text(slide, Inches(0.7), Inches(1.05), Inches(12), Inches(1.0),
             data["title"], font=TITLE_FONT, size=36, bold=True, color=INK_DARK)
    add_text(slide, Inches(0.7), Inches(2.1), Inches(11), Inches(0.6),
             data["lead"], font=SERIF_FONT, size=17, italic=True, color=INK_BODY)

    # four state cards
    states = data["states"]
    n = len(states)
    card_w = 2.8
    card_h = 3.4
    gap = 0.25
    total_w = n * card_w + (n - 1) * gap
    left0 = (13.333 - total_w) / 2
    y0 = 3.2

    for i, (label, desc) in enumerate(states):
        x = left0 + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(x), Inches(y0),
                                       Inches(card_w), Inches(card_h))
        card.fill.solid(); card.fill.fore_color.rgb = BG_CREAM
        card.line.color.rgb = INK_DARK
        card.line.width = Pt(1.25)
        card.shadow.inherit = False

        # accent header
        head = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(x), Inches(y0),
                                       Inches(card_w), Inches(0.6))
        head.fill.solid(); head.fill.fore_color.rgb = ACCENT
        head.line.fill.background()
        add_text(slide, Inches(x), Inches(y0 + 0.13), Inches(card_w), Inches(0.4),
                 label, font=TITLE_FONT, size=16, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), align="center")

        # big number
        add_text(slide, Inches(x), Inches(y0 + 0.8), Inches(card_w), Inches(1.0),
                 f"{i+1:02d}", font=TITLE_FONT, size=46, bold=True,
                 color=ACCENT_SOFT, align="center")

        add_text(slide, Inches(x + 0.2), Inches(y0 + 1.9), Inches(card_w - 0.4), Inches(1.4),
                 desc, font=BODY_FONT, size=13, color=INK_BODY,
                 align="center", line_spacing=1.3)


RENDERERS = {
    "title":    render_title,
    "section":  render_section,
    "twocol":   render_twocol,
    "dataflow": render_dataflow,
    "states":   render_states,
}


# ============================================================
# Build
# ============================================================
def build(out_path="Project_Presentation.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total = len(SLIDES)

    for i, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        RENDERERS[data["type"]](slide, data)
        if data["type"] not in ("title",):
            slide_number_badge(slide, i, total)

    prs.save(out_path)
    print(f"Wrote {out_path}  ({total} slides)")


if __name__ == "__main__":
    build()
